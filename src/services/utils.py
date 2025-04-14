import json
import re
import os
import pandas as pd
from dotenv import load_dotenv
from supabase import create_client
from openai import OpenAI
from sqlalchemy import create_engine, text
from sqlalchemy.orm import sessionmaker
import io
from src.services.prompts import *
#from prompts import sys_prompt_entity_extraction_prompt,sys_prompt_query_generation_without_location,sys_prompt_query_generation_with_location,user_prompt_sql_query_generation

load_dotenv()

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from reportlab.lib.utils import ImageReader
from reportlab.lib import utils
import boto3
import os
from botocore.exceptions import NoCredentialsError, ClientError
from dotenv import load_dotenv
load_dotenv()

access_key = os.getenv("AWS_ACCESS_KEY")
secret_key = os.getenv("AWS_SECRET_KEY")
# print(access_key, secret_key)
def new_ppt_logic(slides_to_keep,output_file):
    slides_to_keep.append(1)
    input_file = "base.pptx"
    #output_file = "selected_slides.pptx"
    """
    Retain only the specified slides (by slide number) from a PowerPoint file.
    
    Args:
        input_path (str): Path to the input PPTX file.
        output_path (str): Path to save the new PPTX file.
        slides_to_keep (list of int): 1-based indices of slides to keep.
    """
    prs = Presentation(input_file)
    total_slides = len(prs.slides)

    # Convert user-friendly slide numbers (1-based) to zero-based indices
    keep_indices = [n - 1 for n in slides_to_keep if 1 <= n <= total_slides]

    # Remove slides that are not in keep_indices
    # Loop in reverse to avoid messing up the indices as we delete slides
    for idx in reversed(range(total_slides)):
        if idx not in keep_indices:
            # Remove the relationship and then remove the slide from the slide list
            rId = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]

    # Save the modified presentation
    prs.save(output_file)
    print(f"Slides {slides_to_keep} have been retained in {output_file}.")
def create_billboard_presentation(slide_data_list, output_filename):
    """
    Creates a PowerPoint presentation with multiple slides for billboard details.

    Parameters:
        slide_data_list (list of dict): A list of dictionaries containing slide details.
            Each dictionary should have the following keys:
            'location', 'dimension', 'lit_type', 'route', 'google_link', 'image_url'
            Note: 'image_url' should be an S3 path like 's3://bucket_name/path/to/image.png'
        output_filename (str): The name of the output PowerPoint file.
    """
    # Create a PowerPoint presentation
    print("ppt started")
    presentation = Presentation()

    # Slide dimensions (default PowerPoint slide size is 10" x 7.5")
    slide_width = Inches(10)
    slide_height = Inches(7.5)

    # Initialize boto3 client outside the loop for efficiency
    s3 = boto3.client('s3',  aws_access_key_id=access_key, aws_secret_access_key=secret_key)

    for idx, slide_data in enumerate(slide_data_list):
        # Add a new slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # Blank layout

        # Add first textbox with center alignment
        width = Inches(5)
        height = Inches(1)
        left = 2  # Center horizontally
        top = Inches(0.5)  # Set top position
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = f"""        District: {slide_data['district']}
        
        Area: {slide_data['area']}
        
        Hoarding Code: {slide_data['hoarding_code']}
        
        Location: {slide_data['location']}
        
        Dimension: {slide_data['dimension']}
        
        Lighting type: {slide_data['lit_type']}
        
        Route: {slide_data['route']}
        
        Location link: {slide_data['google_link']}
        """
        p.font.bold = True
        p.font.size = Pt(20)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center-align text
    
        # # Add route information
        # # Add second textbox with center alignment
        # width = Inches(5)
        # height = Inches(1)
        # left = (slide_width - width) / 2  # Center horizontally
        # top = Inches(1.0)  # Adjust top position for second textbox
        # textbox = slide.shapes.add_textbox(left, top, width, height)
        # text_frame = textbox.text_frame
        # p = text_frame.add_paragraph()
        # p.text = f"{slide_data['route']}"
        # p.font.bold = True
        # p.font.size = Pt(26)
        # text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center-align text
        
        width = Inches(7)
        height = Inches(1)
        left = (slide_width - width) / 2  # Center horizontally
        top = Inches(6.2)  # Adjust top position for second textbox
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = f"operations@palettem.com || Ph no: 9562058880/ 9388140645 || www.palettecommunications.in"
        p.font.bold = True
        p.font.size = Pt(12)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # width = Inches(5)
        # height = Inches(1)
        # left = (slide_width - width) / 2  # Center horizontally
        # top = Inches(6.6)  # Adjust top position for second textbox
        # textbox = slide.shapes.add_textbox(left, top, width, height)
        # text_frame = textbox.text_frame
        # p = text_frame.add_paragraph()
        # p.text = f"Location link: {slide_data['google_link']}"
        # p.font.bold = True
        # p.font.size = Pt(12)
        # text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Add image from S3
        image_path = f"local_image_{idx}.png"  # Unique image path for each slide
        s3_path = slide_data['image_url']  # Expected format: 's3://bucket_name/path/to/image.png'

        # Parse the S3 URL
        if s3_path.startswith("s3://"):
            s3_path = s3_path[5:]  # Remove 's3://'

        try:
            bucket_name, key = s3_path.split('/', 1)  # Split into bucket name and key

            # Download the image from S3
            s3.download_file(bucket_name, key, image_path)
            print(f"Image downloaded from S3 bucket '{bucket_name}' with key '{key}'.")
        except ValueError:
            print("Error: Invalid S3 path format. Please use 's3://bucket_name/path/to/image.png'")
            continue  # Skip to next slide
        except NoCredentialsError:
            print("Error: AWS credentials not found. Please configure your AWS credentials.")
            return
        except ClientError as e:
            print(f"Error downloading image from S3: {e}")
            continue  # Skip to next slide

        # Add image to slide
        width = Inches(4.5)  # Image width
        height = Inches(3.5)  # Image height
        left = (slide_width - width) - 2  # Center horizontally
        top = Inches(1.0)  # Vertical position
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)


    # Save the presentation
    presentation.save(output_filename)
    print(f"Presentation saved as '{output_filename}'.")
    
    

def create_document_with_images(page_data_list, output_filename):
    """
    Creates a Word document with multiple pages, each containing text and an image.

    Parameters:
        page_data_list (list of dict): A list of dictionaries containing page details.
            Each dictionary should have the following keys:
            'location', 'dimension', 'lit_type', 'route', 'google_link', 'image_url'
            'image_url' should be an S3 path like 's3://bucket_name/path/to/image.png'
        output_filename (str): The name of the output Word document file.
    """
    # Create a new Word document
    document = Document()
    
    # Set document style (optional)
    style = document.styles['Normal']
    font = style.font
    font.name = 'Arial'
    font.size = Pt(12)
    
    # Initialize boto3 client for S3
    s3 = boto3.client('s3', aws_access_key_id=access_key, aws_secret_access_key=secret_key)

    for idx, page_data in enumerate(page_data_list):
        # Add a heading
        # title_text = f"{page_data['location']} || {page_data['dimension']} || {page_data['lit_type']}"
        # title = document.add_heading(title_text, level=1)
        # title.alignment = WD_ALIGN_PARAGRAPH.CENTER  # Center-align the title
        
        # Add route information
        route_paragraph = document.add_paragraph()
        route_run = route_paragraph.add_run(f"{page_data['location']} || {page_data['dimension']} || {page_data['lit_type']}")
        route_run.bold = True
        route_run.font.size = Pt(14)
        route_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Add route information
        route_paragraph = document.add_paragraph()
        route_run = route_paragraph.add_run(page_data['route'])
        route_run.bold = True
        route_run.font.size = Pt(14)
        route_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Add image from S3
        image_path = f"local_image_{idx}.png"  # Unique image path for each page
        s3_path = page_data['image_url']  # Expected format: 's3://bucket_name/path/to/image.png'

        # Parse the S3 URL
        if s3_path.startswith("s3://"):
            s3_path = s3_path[5:]  # Remove 's3://'

        try:
            bucket_name, key = s3_path.split('/', 1)  # Split into bucket name and key

            # Download the image from S3
            s3.download_file(bucket_name, key, image_path)
            print(f"Image downloaded from S3 bucket '{bucket_name}' with key '{key}'.")
        except ValueError:
            print("Error: Invalid S3 path format. Please use 's3://bucket_name/path/to/image.png'")
            continue  # Skip to next page
        except NoCredentialsError:
            print("Error: AWS credentials not found. Please configure your AWS credentials.")
            return
        except ClientError as e:
            print(f"Error downloading image from S3: {e}")
            continue  # Skip to next page

        # Add the image to the document
        try:
            picture = document.add_picture(image_path, width=Inches(6))
            last_paragraph = document.paragraphs[-1]
            last_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        except Exception as e:
            print(f"Error adding image on page {idx+1}: {e}")

        # Add contact information
        contact_text = "operations@palettem.com || Ph no: 9562058880/ 9388140645 || www.palettecommunications.in"
        contact_paragraph = document.add_paragraph(contact_text)
        contact_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

        link_text = f"Location link: {page_data['google_link']}"
        link_paragraph = document.add_paragraph(link_text)
        link_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Add a page break after each page except the last one
        if idx < len(page_data_list) - 1:
            document.add_page_break()

    # Save the document
    document.save(output_filename)
    print(f"Document saved as '{output_filename}'.")
    
def create_pdf_with_images(page_data_list, output_filename):
    """
    Creates a PDF document with multiple pages, each containing text and an image.

    Parameters:
        page_data_list (list of dict): A list of dictionaries containing page details.
            Each dictionary should have the following keys:
            'location', 'dimension', 'lit_type', 'route', 'google_link', 'image_url'
            'image_url' should be an S3 path like 's3://bucket_name/path/to/image.png'
        output_filename (str): The name of the output PDF file.
    """
    # Create a canvas object
    c = canvas.Canvas(output_filename, pagesize=A4)
    width, height = A4

    # Initialize boto3 client for S3
    s3 = boto3.client('s3', aws_access_key_id=access_key, aws_secret_access_key=secret_key)

    for idx, page_data in enumerate(page_data_list):
        # --- Page Content Setup ---
        # Set default font and color
        c.setFillColorRGB(0, 0, 0)  # Black color
        c.setFont("Helvetica-Bold", 20)

        # Add Title: location || dimension || lit_type
        title_text = f"{page_data['location']} || {page_data['dimension']} || {page_data['lit_type']}"
        c.drawCentredString(width / 2.0, height - 1 * inch, title_text)

        # Add Route Information
        c.setFont("Helvetica-Bold", 16)
        route_text = page_data['route']
        c.drawCentredString(width / 2.0, height - 1.5 * inch, route_text)

        # --- Download Image from S3 ---
        image_path = f"local_image_{idx}.png"  # Unique image path for each page
        s3_path = page_data['image_url']  # Expected format: 's3://bucket_name/path/to/image.png'

        # Parse the S3 URL
        if s3_path.startswith("s3://"):
            s3_path = s3_path[5:]  # Remove 's3://'

        try:
            bucket_name, key = s3_path.split('/', 1)  # Split into bucket name and key

            # Download the image from S3
            s3.download_file(bucket_name, key, image_path)
            print(f"Image downloaded from S3 bucket '{bucket_name}' with key '{key}'.")
        except ValueError:
            print("Error: Invalid S3 path format. Please use 's3://bucket_name/path/to/image.png'")
            continue  # Skip to next page
        except NoCredentialsError:
            print("Error: AWS credentials not found. Please configure your AWS credentials.")
            return
        except ClientError as e:
            print(f"Error downloading image from S3: {e}")
            continue  # Skip to next page

        # --- Add Image to PDF ---
        try:
            img = utils.ImageReader(image_path)
            img_width, img_height = img.getSize()
            aspect = img_height / float(img_width)

            # Set desired image width (e.g., 6 inches)
            img_display_width = 6 * inch
            img_display_height = img_display_width * aspect

            # Calculate positions to center the image
            img_x = (width - img_display_width) / 2
            img_y = height / 2 - img_display_height / 2

            c.drawImage(image_path, img_x, img_y, img_display_width, img_display_height)
        except Exception as e:
            print(f"Error adding image on page {idx+1}: {e}")

        # --- Add Contact Information ---
        c.setFont("Helvetica", 12)
        contact_text = "operations@palettem.com || Ph no: 9562058880/ 9388140645 || www.palettecommunications.in"
        c.drawCentredString(width / 2.0, 1 * inch, contact_text)

        # Add Google Link
        link_text = f"Location link: {page_data['google_link']}"
        c.drawCentredString(width / 2.0, 0.75 * inch, link_text)

        # --- Finalize Page ---
        c.showPage()

        # Optionally, delete the local image file to clean up
        if os.path.exists(image_path):
            os.remove(image_path)

    # Save the PDF file
    c.save()
    print(f"PDF saved as '{output_filename}'.")













# Initialize Supabase client and OpenAI
#supabase = create_client(st.secrets["SUPABASE_URL"], st.secrets["SUPABASE_KEY"])
openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# Initialize Google Maps client
#gmaps = googlemaps.Client(key=st.secrets["GOOGLE_MAPS_API_KEY"])

# Create the SQLAlchemy engine
# database_url = os.getenv('DATABASE_URL')
# engine = create_engine(database_url)
# Session = sessionmaker(bind=engine)
# session = Session()

##################################################################
# 1. Helper function to geocode a location using Google Maps API #
##################################################################
# def geocode_location(location_str: str):
#     """
#     Uses Google Maps Geocoding API to get latitude and longitude.
#     Returns (lat, lng) or (None, None) if not found.
#     """
#     try:
#         geocode_result = gmaps.geocode(location_str)
#         if geocode_result:
#             lat = geocode_result[0]['geometry']['location']['lat']
#             lng = geocode_result[0]['geometry']['location']['lng']
#             return lat, lng
#         else:
#             return None, None
#     except:
#         return None, None
import requests
GOOGLE_MAPS_API_KEY = os.getenv("GOOGLE_MAPS_API_KEY")#'GOOGLE_MAPS_API_KEY']
def get_lat_long(location: str):
    url = f"https://maps.googleapis.com/maps/api/geocode/json?address={location}&key={GOOGLE_MAPS_API_KEY}"
    response = requests.get(url)
    if response.status_code == 200:
        data = response.json()
        if data['results']:
            lat_lng = data['results'][0]['geometry']['location']
            lat = lat_lng['lat']
            lng = lat_lng['lng']
            return lat, lng
    return None, None

##################################################
# 2. Function to extract entities from user text #
##################################################
def extract_entities(question: str) -> dict:
    """
    Extract the relevant entities from user input using OpenAI chat completion.
    We add logic for picking up 'location' and 'radius'.
    """
    response = openai_client.chat.completions.create(
        model="gpt-4o",
        messages=[#location_route
            {
                "role": "system",
                "content": sys_prompt_entity_extraction_prompt
            },
            {"role": "user", "content": question}
        ]
    )

    # The assistant's JSON is in response.choices[0].message.content
    # We'll parse it:
    raw_response = response.choices[0].message.content.strip()
    raw_response_clean = raw_response.replace("```json", "").replace("```", "").strip()

    # Convert to Python dict
    entities = json.loads(raw_response_clean)

    # If you want to do additional parse checks or fallback:
    # e.g. if "radius" is not in entities, do entities["radius"] = "none"
    
    return entities

def check_input_query(question: str) -> dict:
    """
    Extract the relevant entities from user input using OpenAI chat completion.
    We add logic for picking up 'location' and 'radius'.
    """
    response = openai_client.chat.completions.create(
        model="gpt-4o",
        messages=[#location_route
            {
                "role": "system",
                "content": """
                You role is to determine if the given question is a generic question, or a question related to hoardings.
                If the query is related to how are you or who are you, reply with the response.
                If the query is related to hoardings, reply with lower case text 'hoarding_question'
                If the query is anything else, reply 'I can't answer this question, please go ahead and ask me questions related to hoardings'.
                
                Example:
                Question: How are you?
                Response: I am doing great! Thanks for asking. Please go ahead and ask me the questions.
                
                Question: Who are you?
                Response: I am an AI bot developed for Palette communications.
                
                Question: What is an apple?
                Response: I can't answer this question, please go ahead and ask me questions related to hoardings
                
                Questions: Give me available hoardings in Ernakulam
                Response: hoarding_question"""
            },
            {"role": "user", "content": question}
        ]
    )

    # The assistant's JSON is in response.choices[0].message.content
    # We'll parse it:
    raw_response = response.choices[0].message.content.strip()

    # Convert to Python dic
    # If you want to do additional parse checks or fallback:
    # e.g. if "radius" is not in entities, do entities["radius"] = "none"
    
    return raw_response







def generate_sql_query_GPT2(entities_json: dict, question: str) -> str:
    """
    Given a dictionary of extracted entities, generate a SQL query with adjusted logic for
    location and radius-based filtering using lat and long as separate columns in the database.
    """
    
    location_str = entities_json.get("location", "none")
    radius_str = entities_json.get("radius", "none")
    
    if location_str != "none":
        if radius_str == "none":
            radius_km = 5
        else:
            try:
                radius_km = float(radius_str)
            except:
                radius_km = 5
        radius_in_meters = radius_km * 1000
        lat, lng = get_lat_long(location_str)  # Function to geocode location to lat, lng
        #st.write(lat,lng,radius_in_meters,location_str)
        print("ENTITIES JSON: ", entities_json)
        system_prompt = sys_prompt_query_generation_with_location
        # print("system_prompt",system_prompt)
        user_prompt = user_prompt_loc_sql_query_generation.format(query=question, entities_json=entities_json, lat=lat, lng=lng, radius_in_meters=radius_in_meters)
    else:
        system_prompt = sys_prompt_query_generation_without_location
        # print("system_prompt",system_prompt)
        user_prompt = user_prompt_no_loc_sql_query_generation.format(query=question, entities_json=entities_json)
    
    response = openai_client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0.2
    )
    
    return response.choices[0].message.content.strip().replace("```sql","").replace("```","").strip()

#################################################################
# 5. Execute the SQL query on Supabase or locally via SQLAlchemy #
#################################################################
# def run_sql_query(sql_query: str):
#     # If you have a Supabase RPC that executes SQL:
#     # response = supabase.rpc("execute_sql", {"query": sql_query}).execute()
#     # return response

#     # OR do it directly with SQLAlchemy:
#     try:
#         result = session.execute(text(sql_query))
#         column_names = result.keys()
#         doc = {col: [] for col in column_names}
#         for row in result:
#             for i, col in enumerate(column_names):
#                 doc[col].append(row[i])
#         df = pd.DataFrame(doc)
#         return df
#     except Exception as e:
#         st.error(str(e))
#         return pd.DataFrame()


#########################################
# 6. Full pipeline to process the query #
#########################################
# def process_query(question: str):
#     with st.spinner('Processing your query...'):
#         try:
#             # Step 1: Extract Entities
#             entities = extract_entities(question)
#             st.write("Extracted Entities:", entities)

#             # Step 2: Generate SQL Query
#             #sql_query = generate_sql_query(entities, question)
#             sql_query =generate_sql_query_GPT2(entities, question)
#             st.write("Generated SQL Query:", sql_query)

#             # Step 3: Run SQL Query
#             data_df = run_sql_query(sql_query)
#         except Exception as e:
#             st.error("Please try another query! " + str(e))
#             data_df = pd.DataFrame()

#         return data_df


###################################
# 7. Finally, the Streamlit app   #
###################################

# import streamlit as st
# from utils import create_billboard_presentation, create_document_with_images, create_pdf_with_images
# import pandas as pd
# import os

# Directory to save the generated files
OUTPUT_DIR = "output"
os.makedirs(OUTPUT_DIR, exist_ok=True)

from datetime import datetime

def generate_and_download_files(data, formats,v):
    output_files = []

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")  # Add timestamp for uniqueness

    if "PPT" in formats:
        ppt_file = os.path.join(OUTPUT_DIR, f"Billboard_Presentation_{timestamp}.pptx")
        #create_billboard_presentation(data, ppt_file)
        new_ppt_logic(v, ppt_file)
        output_files.append(ppt_file)

    if "PDF" in formats:
        pdf_file = os.path.join(OUTPUT_DIR, f"Billboard_Document_{timestamp}.pdf")
        create_pdf_with_images(data, pdf_file)
        output_files.append(pdf_file)

    if "DOCX" in formats:
        docx_file = os.path.join(OUTPUT_DIR, f"Billboard_Document_{timestamp}.docx")
        create_document_with_images(data, docx_file)
        output_files.append(docx_file)

    if "XLSX" in formats:
        xlsx_file = os.path.join(OUTPUT_DIR, f"Billboard_Data_{timestamp}.xlsx")
        df = pd.DataFrame(data)
        df.to_excel(xlsx_file, index=False)
        output_files.append(xlsx_file)

    return output_files

# def main():
#     st.title("Hoarding Query System with Download Options")

#     # Example data to simulate user query results
#     example_data = [
#         {
#             'district': 'Ernakulam',
#             'area': 500,
#             'hoarding_code': '1R',
#             'location': 'Kacheripady',
#             'dimension': '20x15',
#             'lit_type': 'Non-Lit',
#             'route': 'North Overbridge to Kacheripady',
#             'google_link': 'https://maps.app.goo.gl/RVjSdRQ9EHFcdDE36',
#             'image_url': 's3://palette-backend/images/10.png',
#         },
#         {
#             'district': 'Ernakulam',
#             'area': 500,
#             'hoarding_code': '1R',
#             'location': 'MG Road',
#             'dimension': '25x10',
#             'lit_type': 'Lit',
#             'route': 'South Railway Station to MG Road',
#             'google_link': 'https://maps.app.goo.gl/RVjSdRQ9EHFcdDE36',
#             'image_url': 's3://palette-backend/images/100.png',
#         },
#     ]

    

def convert_response_to_example_data(response_df):
    """
    Converts the given DataFrame to the format required by the example_data list.

    Parameters:
        response_df (pd.DataFrame): The DataFrame containing hoarding data.

    Returns:
        list: A list of dictionaries in the required format.
    """
    # Rename and select relevant columns
    prepared_data = response_df.rename(
        columns={
            "district": "district",
            "location_route": "location",
            "width": "width",
            "height": "height",
            "area": "area",
            "type": "lit_type",
            "direction_route": "route",
            "hoarding_code": "hoarding_code",
            "lat": "lat",
            "long": "long",
        }
    )[
        [
            "district",
            "area",
            "hoarding_code",
            "location",
            "width",
            "height",
            "lit_type",
            "route",
            "lat",
            "long",
            "hoarding_id",
        ]
    ]

    # Add google_link as the Google Maps URL based on lat and long
    prepared_data["google_link"] = prepared_data.apply(
        lambda row: f"https://www.google.com/maps?q={row['lat']},{row['long']}",
        axis=1
    )

    # Add dimension as "width x height"
    prepared_data["dimension"] = prepared_data.apply(
        lambda row: f"{row['width']} x {row['height']}", axis=1
    )

    # Add image_url (S3 path) based on hoarding_code
    prepared_data["image_url"] = prepared_data.apply(
        lambda row: f"s3://palette-backend/images/{row['hoarding_code']}.png",
        axis=1
    )

    # Reorder columns so the new columns appear in the final DataFrame
    prepared_data = prepared_data[
        [ "hoarding_id",
            "district",
            "area",
            "hoarding_code",
            "location",
            "width",
            "height",
            "lit_type",
            "route",
            "google_link",
            "lat",
            "long",
            "dimension",
            "image_url",
        ]
    ]

    # Convert to list of dicts
    return prepared_data.to_dict(orient="records")


import os
import io
import zipfile

def zip_entire_directory(directory_path: str) -> io.BytesIO:
    """
    Zips all files (and subdirectories) inside 'directory_path' into a single zip file,
    stored in a BytesIO object. Returns that BytesIO (seeked to 0).
    """
    zip_buffer = io.BytesIO()

    # 'w' mode = create a new zip
    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        # Walk through the directory
        for root, dirs, files in os.walk(directory_path):
            for file_name in files:
                full_path = os.path.join(root, file_name)
                # Compute the archive name so that subfolders remain intact
                # relative to the top-level directory_path
                arcname = os.path.relpath(full_path, start=directory_path)
                # Add the file to the zip
                zf.write(full_path, arcname=arcname)

    # Reset pointer so it can be read from the start
    zip_buffer.seek(0)
    return zip_buffer
import shutil
def cleanup_output_dir(directory_path: str):
    """
    Deletes everything inside the given directory_path.
    """
    if os.path.exists(directory_path):
        shutil.rmtree(directory_path)  # removes the entire directory
        os.makedirs(directory_path)
def simple_gpt(user_prompt: str,system_prompt:str,model:str="gpt-4o-mini") -> str:
    """
    Extract the relevant entities from user input using OpenAI chat completion.
    We add logic for picking up 'location' and 'radius'.
    """
    response = openai_client.chat.completions.create(
        model=model,
        messages=[#location_route
            {
                "role": "system",
                "content": system_prompt
            },
            {"role": "user", "content": user_prompt}
        ]
    )

    # The assistant's JSON is in response.choices[0].message.content
    # We'll parse it:
    raw_response = response.choices[0].message.content.strip()
    raw_response_clean = raw_response#.replace("```json", "").replace("```", "").strip()

    # Convert to Python dict
    #entities = json.loads(raw_response_clean)

    # If you want to do additional parse checks or fallback:
    # e.g. if "radius" is not in entities, do entities["radius"] = "none"
    
    return raw_response_clean
# def main():
#     st.title("Hoarding Query System")

#     # Initialize chat history
#     if "messages" not in st.session_state:
#         st.session_state.messages = []

#     # Display chat messages from history on app rerun
#     for message in st.session_state.messages:
#         with st.chat_message(message["role"]):
#             st.markdown(message["content"])

#     # React to user input
#     if prompt := st.chat_input("What would you like to know about hoardings?"):
#         # Display user message in chat
#         st.chat_message("user").markdown(prompt)
#         st.session_state.messages.append({"role": "user", "content": prompt})

#         # Process the query
#         response_df = process_query(prompt)

#         # Display results
#         with st.chat_message("assistant"):
#             if not response_df.empty:
#                 st.dataframe(response_df)
#             else:
#                 st.write("No results found.")

#         # Add assistant response to chat history
#         # (You might store the DataFrame as a string or do something else)
#         st.session_state.messages.append(
#             {
#                 "role": "assistant",
#                 "content": f"Returned {len(response_df)} rows from the database."
#             }
#         )
#         if st.button("Download Results"):
#             with st.expander("Select file formats"):
#                 formats = st.multiselect("Choose formats to download", ["PPT", "PDF", "DOCX", "XLSX"])

#             if formats:
#                 example_data = convert_response_to_example_data(response_df)
#                 print("example_data",example_data)
#                 st.markdown(example_data)
#                 output_files = generate_and_download_files(example_data, formats)

#                 st.success("Files generated successfully!")
#                 for file_path in output_files:
#                     file_name = os.path.basename(file_path)
#                     with open(file_path, "rb") as file:
#                         st.download_button(
#                             label=f"Download {file_name}",
#                             data=file,
#                             file_name=file_name,
#                             mime="application/octet-stream"
#                         )
#             else:
#                 st.warning("Please select at least one format.")
# def main():
#     st.title("Hoarding Query System")

#     # Initialize chat history
#     if "messages" not in st.session_state:
#         st.session_state.messages = []

#     # Display chat messages from history on app rerun
#     for message in st.session_state.messages:
#         with st.chat_message(message["role"]):
#             st.markdown(message["content"])

#     # React to user input
#     if prompt := st.chat_input("What would you like to know about hoardings?"):
#         # Display user message in chat
#         st.chat_message("user").markdown(prompt)
#         st.session_state.messages.append({"role": "user", "content": prompt})

#         # Process the query
#         response_df = process_query(prompt)

#         # Display results
#         with st.chat_message("assistant"):
#             if not response_df.empty:
#                 st.dataframe(response_df)
#             else:
#                 st.write("No results found.")

#         # Add assistant response to chat history
#         st.session_state.messages.append(
#             {
#                 "role": "assistant",
#                 "content": f"Returned {len(response_df)} rows from the database."
#             }
#         )

#         # --- Download Section ---
#         # Prompt the user to select which file formats they want:
#         # with st.expander("Select file formats for download"):
#         #     selected_formats = st.multiselect(
#         #         "Choose one or more formats:",
#         #         ["PPT", "PDF", "DOCX", "XLSX"]
#         #     )
#         # with st.spinner("Getting data"):
#         # # Button to trigger file generation and download
#         #     if st.button("Download Results"):
#         #         if not selected_formats:
#         #             st.warning("Please select at least one format.")
#         #         else:
#         # Convert DF to example_data
#         example_data = convert_response_to_example_data(response_df)
#         st.markdown(example_data)
#         # Generate desired files
#         output_files = generate_and_download_files(example_data, ["PPT", "PDF", "DOCX", "XLSX"])

#         st.success("Files generated successfully!")
#         # Provide a download button for each generated file
#         for file_path in output_files:
#             file_name = os.path.basename(file_path)
#             with open(file_path, "rb") as file:
#                 st.download_button(
#                     label=f"Download {file_name}",
#                     data=file.read(),
#                     file_name=file_name,
#                     mime="application/octet-stream"
#                 )

# if __name__ == "__main__":
#     main()


def generate_text_content(query, chunk_data):
    user_prompt = "Question: " + query + " Available context:\n" + "\n".join(chunk_data)
    return simple_gpt(
        user_prompt=user_prompt,
        system_prompt=system_prompt_text,
        model="gpt-4o"
    )

def generate_blog_content(query, chunk_data):
    user_prompt = "Question: " + query + " Available context:\n" + "\n".join(chunk_data)    
    return simple_gpt(
        user_prompt=user_prompt,
        system_prompt=system_prompt_blog,
        model="gpt-4o"
    )

def generate_ppt_content(query, chunk_data):
    user_prompt = "Question: " + query + " Available context:\n" + "\n".join(chunk_data)
    return simple_gpt(
        user_prompt=user_prompt,
        system_prompt=system_prompt_ppt,
        model="gpt-4o"
    )

def generate_image_content(query, chunk_data):
    user_prompt = "Question: " + query + " Available context:\n" + "\n".join(chunk_data)
    system_prompt = """
You are generating alt-text/image description from context. Return JSON like:
{"image_description": "Detailed, creative description of the image content based on the query and context"}
Return only the JSON.
"""
    
    return simple_gpt(
        user_prompt=user_prompt,
        system_prompt=system_prompt,
        model="gpt-4o"
    )
