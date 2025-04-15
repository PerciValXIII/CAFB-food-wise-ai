import os
import json
import pandas as pd
from dotenv import load_dotenv
from pptx import Presentation
import boto3
from supabase import create_client, Client

#loading env variables
load_dotenv()

# Supabase client setup
SUPABASE_URL = os.getenv("SUPABASE_URL")   
#service role secret key (not anon public)
SUPABASE_KEY = os.getenv("SUPABASE_KEY")  
TABLE_NAME = "ppt_data"
aws_access_key_id = os.getenv("AWS_ACCESS_KEY_ID")
aws_secret_access_key = os.getenv("AWS_SECRET_ACCESS_KEY")
aws_default_region = os.getenv("AWS_DEFAULT_REGION")
root_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ppt_dir = os.path.join(root_dir, "data", "sources", "powerpoints")
output_csv = os.path.join(root_dir, "data", "supabase_structured_data", "extracted_ppts", "parsed_ppt_slides.csv")
bucket_name = "cfab"
s3_folder = "on_premise_data/powerpoints/"
log_file = os.path.join(root_dir, "uploaded_files.txt")


#Initializing Supabase client
supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY)

#Initialize s3 client
s3 = boto3.client(
    's3',
    aws_access_key_id=aws_access_key_id,
    aws_secret_access_key=aws_secret_access_key,
    region_name=aws_default_region  
)

#uploading user files to s3
for filename in os.listdir(ppt_dir):
    local_path = os.path.join(ppt_dir, filename)
    
    if os.path.isfile(local_path):
        s3_path = os.path.join(s3_folder, filename).replace("\\", "/")  
        print(f"Uploading {filename} to s3://{bucket_name}/{s3_path}")
        
        s3.upload_file(local_path, bucket_name, s3_path)

print("Upload complete!")

#checking supabase table records
response = supabase.table(TABLE_NAME).select("file_name").execute()
supabase_file_names = {entry["file_name"] for entry in response.data}
print(f"Files currently in Supabase: {len(supabase_file_names)}")

ppt_files = [f for f in os.listdir(ppt_dir) if f.endswith(".pptx")]
new_files = [f for f in ppt_files if f not in supabase_file_names]
print(f"Files to process: {new_files}")
print(f"Found {len(new_files)} new PowerPoint file(s) to process and upload.")

#file_id
all_file_ids = supabase.table(TABLE_NAME).select("file_id").execute()
existing_ids = [int(entry["file_id"].split("_")[1]) for entry in all_file_ids.data if "file_id" in entry]
starting_index = max(existing_ids) + 1 if existing_ids else 1
print(f"Starting from file_id: ppt_{starting_index:03d}")

class PowerPointParser:
    def __init__(self, ppt_path):
        self.ppt_path = ppt_path
        self.file_name = os.path.basename(ppt_path)
        self.presentation = Presentation(ppt_path)
        self.file_type = "powerpoint"

    def extract_text_from_slide(self, slide, slide_number):
        slide_title = ""
        content_blocks = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if not para_text:
                    continue

                if not slide_title:
                    slide_title = para_text
                    continue

                if paragraph.level > 0 or para_text.startswith("•") or para_text.startswith("-"):
                    para_type = "bullet"
                else:
                    para_type = "paragraph"

                content_blocks.append({
                    "type": para_type,
                    "text": para_text
                })

        structured_json = {
            "filename": self.file_name,
            "slide_number": slide_number,
            "slide_title": slide_title,
            "content": content_blocks
        }
        return structured_json

    def parse_presentation(self, file_id):
        data = []
        for idx, slide in enumerate(self.presentation.slides, start=1):
            slide_json = self.extract_text_from_slide(slide, idx)
            data.append({
                "file_id": file_id,
                "file_name": self.file_name,
                "slide_number": idx,
                "content_json": json.dumps(slide_json),
                "file_type": self.file_type,
                "train": False
            })
        return data


for i, ppt_file in enumerate(sorted(new_files), start=starting_index):
    file_id = f"ppt_{i:03d}"
    ppt_path = os.path.join(ppt_dir, ppt_file)
    print(f"Processing {ppt_file} as {file_id}...")

    parser = PowerPointParser(ppt_path)
    parsed_data = parser.parse_presentation(file_id)

    if parsed_data:
        try:
            response = supabase.table(TABLE_NAME).insert(parsed_data).execute()
            print(f"Uploaded {ppt_file}: {response}")

            # Log it locally if needed
            with open(log_file, "a") as f:
                f.write(ppt_file + "\n")
        except Exception as e:
            print(f"Error uploading {ppt_file}: {e}")
    else:
        print(f"No data extracted from {ppt_file}. Skipping.")

print("All files processed.")