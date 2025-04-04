"""
================================================================================
Script Name: sm_ppt_structurization.py
Author: Swattik Maiti
Description:
    This script processes all PowerPoint (.pptx) files located in:
        mainfolder/data/sources/powerpoints/

    For each presentation, it:
    - Iterates through each slide
    - Extracts the title and content from slide shapes
    - Differentiates between 'paragraph' and 'bullet' text using indentation levels
    - Structures each slide's content into a JSON object with the following keys:
        - filename
        - slide_number
        - slide_title
        - content: list of text blocks with 'type' (paragraph/bullet) and 'text'

    Each slide is represented as a row in the final CSV:
        mainfolder/data/supabase_structured_data/parsed_ppt_slides.csv

    The CSV contains the following columns:
        - file_id      (e.g., ppt_001)
        - file_name    (original .pptx filename)
        - slide_number
        - content_json (JSON structure of slide content)
        - file_type    (always set to "powerpoint")

Dependencies:
    - python-pptx
    - pandas
================================================================================
"""


import os
import json
import pandas as pd
from pptx import Presentation


class PowerPointParser:
    def __init__(self, ppt_path):
        self.ppt_path = ppt_path
        self.file_name = os.path.basename(ppt_path)
        self.presentation = Presentation(ppt_path)
        self.file_type = "powerpoint"

    def extract_text_from_slide(self, slide, slide_number):
        slide_title = ""
        content_blocks = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if not para_text:
                    continue

                # Determine paragraph type
                if not slide_title:
                    slide_title = para_text
                    continue

                if paragraph.level > 0 or para_text.startswith("•") or para_text.startswith("-"):
                    para_type = "bullet"
                else:
                    para_type = "paragraph"

                content_blocks.append({
                    "type": para_type,
                    "text": para_text
                })

        structured_json = {
            "filename": self.file_name,
            "slide_number": slide_number,
            "slide_title": slide_title,
            "content": content_blocks
        }
        return structured_json

    def parse_presentation(self, file_id):
        data = []
        for idx, slide in enumerate(self.presentation.slides, start=1):
            slide_json = self.extract_text_from_slide(slide, idx)
            data.append({
                "file_id": file_id,
                "file_name": self.file_name,
                "slide_number": idx,
                "content_json": json.dumps(slide_json),
                "file_type": self.file_type
            })
        return data




if __name__ == "__main__":
    root_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    ppt_dir = os.path.join(root_dir, "data", "sources", "powerpoints")
    output_csv = os.path.join(root_dir, "data", "supabase_structured_data","parsed_ppt_slides.csv")

    all_records = []
    ppt_files = [f for f in os.listdir(ppt_dir) if f.endswith(".pptx")]

    for i, ppt_file in enumerate(sorted(ppt_files), start=1):
        file_id = f"ppt_{i:03d}"
        ppt_path = os.path.join(ppt_dir, ppt_file)
        parser = PowerPointParser(ppt_path)
        parsed_data = parser.parse_presentation(file_id)
        all_records.extend(parsed_data)

    df = pd.DataFrame(all_records)
    df.to_csv(output_csv, index=False)
    print(f"Parsed {len(ppt_files)} files. Data saved to: {output_csv}")
