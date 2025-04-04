import io
import os
import json
import datetime
from pathlib import Path
import jsonschema
from pptx import Presentation
from pptx.dml.color import RGBColor
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from PIL import Image

class ScrapePowerPoints:
    def __init__(self, source_dir, output_file, image_save_path="ppt_images", stats_file="ppt_stats.json"):
        """Initialize scraper with source directory and output file path."""
        self.source_dir = Path(source_dir).resolve()
        self.output_file = Path(output_file).resolve()
        self.image_save_path = Path(image_save_path).resolve()
        self.stats_file = Path(stats_file).resolve()
        self.image_metadata_file = self.image_save_path / "powerpoints.jsonl"

        # Ensure output directories exist
        self.output_file.parent.mkdir(parents=True, exist_ok=True)
        self.image_save_path.mkdir(parents=True, exist_ok=True)

        print(f"[INFO] Initialized ScrapePowerPoints for directory: {self.source_dir}")

    def extract_metadata(self, pptx_path):
        """Extract metadata from PowerPoint file."""
        try:
            prs = Presentation(pptx_path)
            core_props = prs.core_properties
            return {
                "Title": core_props.title or "Unknown",
                "Author": core_props.author or "Unknown",
                "Subject": core_props.subject or "Unknown",
                "Revision": core_props.revision or "Unknown",
                "Created": core_props.created.isoformat() if core_props.created else "Unknown",
                "Modified": core_props.modified.isoformat() if core_props.modified else "Unknown",
            }
        except Exception as e:
            print(f"[ERROR] Failed to extract metadata for {pptx_path}: {e}")
            return {}

    def extract_text_with_location(self, pptx_path):
        """Extract text along with slide and paragraph location."""
        slides_data = []
        word_count = 0

        try:
            prs = Presentation(pptx_path)
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for para_index, shape in enumerate(slide.shapes, start=1):
                    if hasattr(shape, "text") and shape.text.strip():
                        words = shape.text.strip().split()
                        word_count += len(words)
                        slide_text.append({
                            "slide": slide_num,
                            "paragraph": para_index,
                            "text": shape.text.strip()
                        })
                slides_data.extend(slide_text)

        except Exception as e:
            print(f"[ERROR] Failed to extract text from {pptx_path}: {e}")

        return slides_data, word_count

    def extract_images(self, pptx_path, ppt_index):
        """Extract images from PowerPoint and save them as PNGs."""
        images = []
        image_metadata = []

        try:
            prs = Presentation(pptx_path)
            for slide_num, slide in enumerate(prs.slides, start=1):
                for img_index, shape in enumerate(slide.shapes, start=1):
                    if hasattr(shape, "image"):
                        img = shape.image

                        # Convert extracted binary image to PNG format
                        image_data = Image.open(io.BytesIO(img.blob))
                        image_filename = f"pow_{ppt_index}_{str(slide_num).zfill(2)}_{str(img_index).zfill(2)}.png"
                        image_path = self.image_save_path / image_filename

                        # Save image as PNG
                        image_data.save(image_path, format="PNG")

                        images.append(str(image_path))

                        # Collect metadata for the extracted image
                        image_metadata.append({
                            "image_name": image_filename,
                            "slide_number": slide_num,
                            "original_ppt": pptx_path.name
                        })

        except Exception as e:
            print(f"[ERROR] Failed to extract images from {pptx_path}: {e}")

        # Save image metadata to JSONL file
        self.save_to_jsonl(image_metadata, self.image_metadata_file, append=True)

        return images

    def validate_jsonl_entry(self, entry):
        """Validate JSON structure before writing."""
        schema = {
            "type": "object",
            "properties": {
                "file_name": {"type": "string"},
                "file_path": {"type": "string"},
                "metadata": {"type": "object"},
                "text_data": {"type": "array"},
                "images": {"type": "array"}
            },
            "required": ["file_name", "file_path", "metadata", "text_data", "images"]
        }
        try:
            jsonschema.validate(instance=entry, schema=schema)
            return True
        except jsonschema.exceptions.ValidationError as e:
            print(f"[ERROR] JSON Validation Failed: {e}")
            return False
    
    def process_ppts(self):
        """Process all PowerPoint files in the source directory and save structured data."""
        ppt_files = sorted(self.source_dir.glob("*.pptx"), key=lambda x: x.name)  # Sort alphanumerically
        all_data = []
        stats_data = []

        if not ppt_files:
            print("[WARNING] No PowerPoint files found in the source directory.")
            return

        print(f"[INFO] Found {len(ppt_files)} PowerPoint files. Processing...")

        for file_index, ppt_file in enumerate(ppt_files, start=1):
            print(f"[INFO] Processing: {ppt_file.name}")

            ppt_index = ppt_file.stem.split("_")[0]  # Extract numeric ID (e.g., 001 from '001_presentation.pptx')

            metadata = self.extract_metadata(ppt_file)
            text_data, word_count = self.extract_text_with_location(ppt_file)
            images = self.extract_images(ppt_file, ppt_index)
            file_link = f"./{ppt_file.relative_to(self.source_dir.parent)}"

            ppt_entry = {
                "file_name": ppt_file.name,
                "file_path": file_link,
                "metadata": metadata,
                "text_data": text_data,
                "images": images
            }

            # Validate JSON before appending
            if self.validate_jsonl_entry(ppt_entry):
                all_data.append(ppt_entry)

            # Save document statistics
            stats_data.append({
                "file_name": ppt_file.name,
                "word_count": word_count,
                "image_count": len(images),
            })

        # Sort final data alphabetically by filename before saving
        all_data.sort(key=lambda x: x["file_name"])

        # Save data to JSONL file
        self.save_to_jsonl(all_data, self.output_file)

        # Save statistics file
        self.save_to_json(stats_data, self.stats_file)

        # Run similarity analysis
        self.compare_text_similarity(all_data)

    def save_to_jsonl(self, data, file_path, append=False):
        """Ensures JSONL is saved properly with one valid JSON object per line."""
        try:
            mode = "a" if append else "w"
            with open(file_path, mode, encoding="utf-8") as file:
                for entry in data:
                    json_str = json.dumps(entry, ensure_ascii=False)  # Convert to string
                    file.write(json_str + "\n")  # Ensure newline to separate JSON objects

            print(f"[INFO] Successfully saved {len(data)} records to {file_path} in JSONL format.")
        except Exception as e:
            print(f"[ERROR] Failed to save JSONL: {e}")



    def save_to_json(self, data, file_path):
        """Saves data to a JSON file."""
        try:
            with open(file_path, "w", encoding="utf-8") as file:
                json.dump(data, file, ensure_ascii=False, indent=4)
            print(f"[INFO] Saved statistics to {file_path}")
        except Exception as e:
            print(f"[ERROR] Failed to save statistics: {e}")

    def compare_text_similarity(self, all_data):
        """Compares the textual content of PPTs using cosine similarity."""
        print("\n[INFO] Running similarity analysis...")

        file_names = [entry["file_name"] for entry in all_data]
        text_contents = [" ".join([t["text"] for t in entry["text_data"]]) for entry in all_data]

        if len(text_contents) < 2:
            print("[WARNING] Not enough documents for similarity comparison.")
            return

        # Convert text into TF-IDF matrix
        vectorizer = TfidfVectorizer(stop_words="english")
        tfidf_matrix = vectorizer.fit_transform(text_contents)

        # Compute cosine similarity
        similarity_matrix = cosine_similarity(tfidf_matrix)

        # Save similarity results
        similarity_data = []
        for i, filename1 in enumerate(file_names):
            for j, filename2 in enumerate(file_names):
                if i < j:
                    similarity_data.append({
                        "file1": filename1,
                        "file2": filename2,
                        "similarity": round(float(similarity_matrix[i, j]), 2)
                    })

        self.save_to_json(similarity_data, self.stats_file)

# Example Usage
if __name__ == "__main__":
    base_path = Path(__file__).resolve().parent
    source_directory = base_path / "../sources"
    output_jsonl = base_path / "powerpoints.jsonl"
    image_dir = base_path / "ppt_images"
    stats_file = base_path / "ppt_stats.json"

    scraper = ScrapePowerPoints(source_directory, output_jsonl, image_dir, stats_file)
    scraper.process_ppts()
