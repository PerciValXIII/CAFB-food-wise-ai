import os
import pandas as pd
from pptx import Presentation
import fitz  # PyMuPDF

def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        doc = fitz.open(pdf_path)
        for page in doc:
            text += page.get_text() + "\n"
        doc.close()
    except Exception as e:
        print(f"Error reading PDF: {pdf_path} | {e}")
    return text.strip()

def extract_text_from_pptx(pptx_path):
    text = ""
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPTX: {pptx_path} | {e}")
    return text.strip()

def format_content_preserving_paragraphs(raw_text):
    return raw_text.replace("\n", " <|endoftext|> ").replace("\r", "").strip()

def fill_file_content(
    csv_relative_path="../data/annotated_data/text_content_annotation.csv",
    ppt_dir_relative="../data/sources/powerpoints/",
    pdf_dir_relative="../data/sources/collateral/"
):
    base_dir = os.path.dirname(os.path.abspath(__file__))
    csv_path = os.path.join(base_dir, csv_relative_path)
    ppt_dir = os.path.join(base_dir, ppt_dir_relative)
    pdf_dir = os.path.join(base_dir, pdf_dir_relative)

    df = pd.read_csv(csv_path)
    content_list = []

    for _, row in df.iterrows():
        file_name = row['File Name']
        source_type = str(row['Source Type']).strip().lower()

        if not isinstance(file_name, str) or not file_name.strip():
            content_list.append("")
            continue

        if 'powerpoint' in source_type:
            file_path = os.path.join(ppt_dir, file_name)
            raw_text = extract_text_from_pptx(file_path)
        elif 'pdf' in source_type or 'collateral' in source_type:
            file_path = os.path.join(pdf_dir, file_name)
            raw_text = extract_text_from_pdf(file_path)
        else:
            raw_text = ""

        formatted_text = format_content_preserving_paragraphs(raw_text)
        content_list.append(formatted_text)

    df['Content'] = content_list

    # Save to CSV
    df.to_csv(csv_path, index=False)

    # Save to Excel (for better viewing)
    excel_path = csv_path.replace(".csv", ".xlsx")
    df.to_excel(excel_path, index=False)

    print("Content column updated with <|endoftext|> formatting and saved to CSV & Excel.")

# Optional: run directly
if __name__ == "__main__":
    fill_file_content()
